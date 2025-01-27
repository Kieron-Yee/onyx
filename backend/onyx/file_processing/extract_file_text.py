"""
文件功能说明：
此模块用于处理和提取各种文件格式的文本内容。支持的文件格式包括：
- 纯文本文件（txt、md、log等）
- PDF文件
- Word文档（docx）
- PowerPoint文件（pptx）
- Excel文件（xlsx）
- 电子邮件（eml）
- 电子书（epub）
- HTML文件
"""

import io
import json
import os
import re
import zipfile
from collections.abc import Callable
from collections.abc import Iterator
from email.parser import Parser as EmailParser
from io import BytesIO
from pathlib import Path
from typing import Any
from typing import Dict
from typing import IO

import chardet
import docx  # type: ignore
import openpyxl  # type: ignore
import pptx  # type: ignore
from docx import Document
from fastapi import UploadFile
from pypdf import PdfReader
from pypdf.errors import PdfStreamError

from onyx.configs.constants import DANSWER_METADATA_FILENAME
from onyx.configs.constants import FileOrigin
from onyx.file_processing.html_utils import parse_html_page_basic
from onyx.file_processing.unstructured import get_unstructured_api_key
from onyx.file_processing.unstructured import unstructured_to_text
from onyx.file_store.file_store import FileStore
from onyx.utils.logger import setup_logger

logger = setup_logger()

"""常量说明"""
# 文本段落之间的分隔符
TEXT_SECTION_SEPARATOR = "\n\n"

# 支持的纯文本文件扩展名列表
PLAIN_TEXT_FILE_EXTENSIONS = [
    ".txt",
    ".md",
    ".mdx",
    ".conf",
    ".log",
    ".json",
    ".csv",
    ".tsv",
    ".xml",
    ".yml",
    ".yaml",
]

# 所有支持的文件扩展名列表
VALID_FILE_EXTENSIONS = PLAIN_TEXT_FILE_EXTENSIONS + [
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".eml",
    ".epub",
    ".html",
]

def is_text_file_extension(file_name: str) -> bool:
    """
    检查文件是否为纯文本文件（根据扩展名判断）
    
    参数:
        file_name: 文件名
    返回:
        bool: 如果是支持的纯文本文件扩展名返回True，否则返回False
    """
    return any(file_name.endswith(ext) for ext in PLAIN_TEXT_FILE_EXTENSIONS)

def get_file_ext(file_path_or_name: str | Path) -> str:
    """
    获取文件的扩展名（小写形式）
    
    参数:
        file_path_or_name: 文件路径或文件名
    返回:
        str: 文件扩展名（小写）
    """
    _, extension = os.path.splitext(file_path_or_name)
    return extension.lower()

def is_valid_file_ext(ext: str) -> bool:
    """
    检查文件扩展名是否在支持的列表中
    
    参数:
        ext: 文件扩展名
    返回:
        bool: 如果扩展名在支持列表中返回True，否则返回False
    """
    return ext in VALID_FILE_EXTENSIONS

def is_text_file(file: IO[bytes]) -> bool:
    """
    通过检查文件内容判断是否为文本文件
    
    检查文件前1024字节是否只包含可打印字符或空白字符。
    如果是，则认为是纯文本文件。
    
    参数:
        file: 文件对象
    返回:
        bool: 如果是文本文件返回True，否则返回False
    """
    raw_data = file.read(1024)
    text_chars = bytearray({7, 8, 9, 10, 12, 13, 27} | set(range(0x20, 0x100)) - {0x7F})
    return all(c in text_chars for c in raw_data)

def detect_encoding(file: IO[bytes]) -> str:
    """
    检测文件的字符编码
    
    参数:
        file: 文件对象
    返回:
        str: 检测到的编码，如果检测失败则返回'utf-8'
    """
    raw_data = file.read(50000)
    encoding = chardet.detect(raw_data)["encoding"] or "utf-8"
    file.seek(0)
    return encoding

def is_macos_resource_fork_file(file_name: str) -> bool:
    """
    检查文件是否为MacOS资源文件
    
    参数:
        file_name: 文件名
    返回:
        bool: 如果是MacOS资源文件返回True,否则返回False
    """
    return os.path.basename(file_name).startswith("._") and file_name.startswith(
        "__MACOSX"
    )

# To include additional metadata in the search index, add a .onyx_metadata.json file
# to the zip file. This file should contain a list of objects with the following format:
# [{ "filename": "file1.txt", "link": "https://example.com/file1.txt" }]
def load_files_from_zip(
    zip_file_io: IO,
    ignore_macos_resource_fork_files: bool = True,
    ignore_dirs: bool = True,
) -> Iterator[tuple[zipfile.ZipInfo, IO[Any], dict[str, Any]]]:
    """
    从ZIP文件中加载文件并解析元数据
    
    要在搜索索引中包含额外的元数据,请在zip文件中添加.onyx_metadata.json文件。
    此文件应包含具有以下格式的对象列表:
    [{ "filename": "file1.txt", "link": "https://example.com/file1.txt" }]
    
    参数:
        zip_file_io: ZIP文件对象
        ignore_macos_resource_fork_files: 是否忽略MacOS资源文件
        ignore_dirs: 是否忽略目录
    返回:
        Iterator[tuple]: 包含(文件信息,文件对象,元数据)的迭代器
    """
    with zipfile.ZipFile(zip_file_io, "r") as zip_file:
        zip_metadata = {}
        try:
            metadata_file_info = zip_file.getinfo(DANSWER_METADATA_FILENAME)
            with zip_file.open(metadata_file_info, "r") as metadata_file:
                try:
                    zip_metadata = json.load(metadata_file)
                    if isinstance(zip_metadata, list):
                        # convert list of dicts to dict of dicts
                        zip_metadata = {d["filename"]: d for d in zip_metadata}
                except json.JSONDecodeError:
                    logger.warn(f"Unable to load {DANSWER_METADATA_FILENAME}")
        except KeyError:
            logger.info(f"No {DANSWER_METADATA_FILENAME} file")

        for file_info in zip_file.infolist():
            with zip_file.open(file_info.filename, "r") as file:
                if ignore_dirs and file_info.is_dir():
                    continue

                if (
                    ignore_macos_resource_fork_files
                    and is_macos_resource_fork_file(file_info.filename)
                ) or file_info.filename == DANSWER_METADATA_FILENAME:
                    continue
                yield file_info, file, zip_metadata.get(file_info.filename, {})

def _extract_onyx_metadata(line: str) -> dict | None:
    """
    从文本行中提取onyx元数据
    
    支持两种格式:
    1. HTML注释格式: <!-- DANSWER_METADATA={...} -->
    2. 标签格式: #DANSWER_METADATA={...}
    
    参数:
        line: 要解析的文本行
    返回:
        dict|None: 解析出的元数据字典,解析失败返回None
    """
    html_comment_pattern = r"<!--\s*DANSWER_METADATA=\{(.*?)\}\s*-->"
    hashtag_pattern = r"#DANSWER_METADATA=\{(.*?)\}"

    html_comment_match = re.search(html_comment_pattern, line)
    hashtag_match = re.search(hashtag_pattern, line)

    if html_comment_match:
        json_str = html_comment_match.group(1)
    elif hashtag_match:
        json_str = hashtag_match.group(1)
    else:
        return None

    try:
        return json.loads("{" + json_str + "}")
    except json.JSONDecodeError:
        return None

def read_text_file(
    file: IO,
    encoding: str = "utf-8",
    errors: str = "replace",
    ignore_onyx_metadata: bool = True,
) -> tuple[str, dict]:
    """
    读取文本文件内容和元数据
    
    参数:
        file: 文件对象
        encoding: 文件编码
        errors: 编码错误处理方式
        ignore_onyx_metadata: 是否忽略onyx元数据
    返回:
        tuple[str, dict]: (文件内容,元数据)元组
    """
    metadata = {}
    file_content_raw = ""
    for ind, line in enumerate(file):
        try:
            line = line.decode(encoding) if isinstance(line, bytes) else line
        except UnicodeDecodeError:
            line = (
                line.decode(encoding, errors=errors)
                if isinstance(line, bytes)
                else line
            )

        if ind == 0:
            metadata_or_none = (
                None if ignore_onyx_metadata else _extract_onyx_metadata(line)
            )
            if metadata_or_none is not None:
                metadata = metadata_or_none
            else:
                file_content_raw += line
        else:
            file_content_raw += line

    return file_content_raw, metadata

def pdf_to_text(file: IO[Any], pdf_pass: str | None = None) -> str:
    """
    从PDF文件中提取文本内容
    
    参数:
        file: PDF文件对象
        pdf_pass: PDF密码(如果有的话)
    返回:
        str: 提取的文本内容
    """
    # Return only the extracted text from read_pdf_file
    text, _ = read_pdf_file(file, pdf_pass)
    return text

def read_pdf_file(
    file: IO[Any],
    pdf_pass: str | None = None,
) -> tuple[str, dict]:
    """
    读取PDF文件并提取文本内容和元数据
    
    参数:
        file: PDF文件对象
        pdf_pass: PDF密码(如果需要)
    返回:
        tuple[str, dict]: 包含提取的文本内容和元数据的元组
    """
    metadata: Dict[str, Any] = {}
    try:
        pdf_reader = PdfReader(file)

        # If marked as encrypted and a password is provided, try to decrypt
        # 如果PDF被加密且提供了密码,尝试解密
        if pdf_reader.is_encrypted and pdf_pass is not None:
            decrypt_success = False
            if pdf_pass is not None:
                try:
                    decrypt_success = pdf_reader.decrypt(pdf_pass) != 0
                except Exception:
                    logger.error("Unable to decrypt pdf")

            if not decrypt_success:
                # By user request, keep files that are unreadable just so they
                # can be discoverable by title.
                return "", metadata
        elif pdf_reader.is_encrypted:
            logger.warning("No Password available to decrypt pdf, returning empty")
            return "", metadata

        # Extract metadata from the PDF, removing leading '/' from keys if present
        # This standardizes the metadata keys for consistency
        # 从PDF中提取元数据,删除键前面的'/',这样可以使元数据键的格式保持一致
        metadata = {}
        if pdf_reader.metadata is not None:
            for key, value in pdf_reader.metadata.items():
                clean_key = key.lstrip("/")
                if isinstance(value, str) and value.strip():
                    metadata[clean_key] = value

                elif isinstance(value, list) and all(
                    isinstance(item, str) for item in value
                ):
                    metadata[clean_key] = ", ".join(value)

        return (
            TEXT_SECTION_SEPARATOR.join(
                page.extract_text() for page in pdf_reader.pages
            ),
            metadata,
        )
    except PdfStreamError:
        logger.exception("PDF file is not a valid PDF")
    except Exception:
        logger.exception("Failed to read PDF")

    # File is still discoverable by title
    # but the contents are not included as they cannot be parsed
    # 文件仍可通过标题被发现,但由于无法解析,内容不会被包含
    return "", metadata

def docx_to_text(file: IO[Any]) -> str:
    """
    从Word文档(docx)中提取文本内容
    
    包括:
    - 正文文本
    - 表格内容(简单表格)
    
    参数:
        file: Word文档文件对象
    返回:
        str: 提取的文本内容
    """
    def is_simple_table(table: docx.table.Table) -> bool:
        """
        判断是否为简单表格(无合并单元格和嵌套表格)
        
        参数:
            table: docx表格对象
        返回:
            bool: 如果是简单表格返回True,否则返回False
        """
        for row in table.rows:
            # 检查是否有被省略的单元格
            if row.grid_cols_before > 0 or row.grid_cols_after > 0:
                return False

            # 检查是否有嵌套的表格
            if any(cell.tables for cell in row.cells):
                return False

        return True

    def extract_cell_text(cell: docx.table._Cell) -> str:
        """
        提取单元格中的文本内容
        
        参数:
            cell: 单元格对象
        返回:
            str: 单元格中的文本内容,如果为空则返回'N/A'
        """
        cell_paragraphs = [para.text.strip() for para in cell.paragraphs]
        return " ".join(p for p in cell_paragraphs if p) or "N/A"

    paragraphs = []
    doc = docx.Document(file)
    for item in doc.iter_inner_content():
        if isinstance(item, docx.text.paragraph.Paragraph):
            paragraphs.append(item.text)

        elif isinstance(item, docx.table.Table):
            if not item.rows or not is_simple_table(item):
                continue

            # 每行作为一个新行,用换行符连接
            table_content = "\n".join(
                [
                    ",\t".join(extract_cell_text(cell) for cell in row.cells)
                    for row in item.rows
                ]
            )
            paragraphs.append(table_content)

    # Word文档已经有适当的段落间距
    return "\n".join(paragraphs)

def pptx_to_text(file: IO[Any]) -> str:
    """
    从PowerPoint文件(pptx)中提取文本内容
    
    参数:
        file: PowerPoint文件对象
    返回:
        str: 按幻灯片顺序提取的文本内容
    """
    presentation = pptx.Presentation(file)
    text_content = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        extracted_text = f"\nSlide {slide_number}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"
        text_content.append(extracted_text)
    return TEXT_SECTION_SEPARATOR.join(text_content)

def xlsx_to_text(file: IO[Any]) -> str:
    """
    从Excel文件(xlsx)中提取文本内容
    
    参数:
        file: Excel文件对象
    返回:
        str: 提取的所有工作表的文本内容
    """
    workbook = openpyxl.load_workbook(file, read_only=True)
    text_content = []
    for sheet in workbook.worksheets:
        sheet_string = "\n".join(
            ",".join(map(str, row))
            for row in sheet.iter_rows(min_row=1, values_only=True)
        )
        text_content.append(sheet_string)
    return TEXT_SECTION_SEPARATOR.join(text_content)

def eml_to_text(file: IO[Any]) -> str:
    """
    从电子邮件文件(eml)中提取文本内容
    
    参数:
        file: 邮件文件对象
    返回:
        str: 提取的邮件文本内容
    """
    text_file = io.TextIOWrapper(file, encoding=detect_encoding(file))
    parser = EmailParser()
    message = parser.parse(text_file)
    text_content = []
    for part in message.walk():
        if part.get_content_type().startswith("text/plain"):
            text_content.append(part.get_payload())
    return TEXT_SECTION_SEPARATOR.join(text_content)

def epub_to_text(file: IO[Any]) -> str:
    """
    从电子书文件(epub)中提取文本内容
    
    参数:
        file: EPUB文件对象
    返回:
        str: 提取的电子书文本内容
    """
    with zipfile.ZipFile(file) as epub:
        text_content = []
        for item in epub.infolist():
            if item.filename.endswith(".xhtml") or item.filename.endswith(".html"):
                with epub.open(item) as html_file:
                    text_content.append(parse_html_page_basic(html_file))
        return TEXT_SECTION_SEPARATOR.join(text_content)

def file_io_to_text(file: IO[Any]) -> str:
    """
    从通用文件中提取文本内容
    
    参数:
        file: 文件对象
    返回:
        str: 提取的文本内容
    """
    encoding = detect_encoding(file)
    file_content_raw, _ = read_text_file(file, encoding=encoding)
    return file_content_raw

def extract_file_text(
    file: IO[Any],
    file_name: str,
    break_on_unprocessable: bool = True,
    extension: str | None = None,
) -> str:
    """
    根据文件类型提取文件中的文本内容
    
    参数:
        file: 文件对象
        file_name: 文件名
        break_on_unprocessable: 当处理失败时是否抛出异常
        extension: 文件扩展名(可选)
    返回:
        str: 提取的文本内容
    异常:
        RuntimeError: 当break_on_unprocessable为True且文件处理失败时抛出
    """
    extension_to_function: dict[str, Callable[[IO[Any]], str]] = {
        ".pdf": pdf_to_text,
        ".docx": docx_to_text,
        ".pptx": pptx_to_text,
        ".xlsx": xlsx_to_text,
        ".eml": eml_to_text,
        ".epub": epub_to_text,
        ".html": parse_html_page_basic,
    }

    try:
        if get_unstructured_api_key():
            return unstructured_to_text(file, file_name)

        if file_name or extension:
            if extension is not None:
                final_extension = extension
            elif file_name is not None:
                final_extension = get_file_ext(file_name)

            if is_valid_file_ext(final_extension):
                return extension_to_function.get(final_extension, file_io_to_text)(file)

        # Either the file somehow has no name or the extension is not one that we recognize
        if is_text_file(file):
            return file_io_to_text(file)

        raise ValueError("Unknown file extension and unknown text encoding")

    except Exception as e:
        if break_on_unprocessable:
            raise RuntimeError(
                f"Failed to process file {file_name or 'Unknown'}: {str(e)}"
            ) from e
        logger.warning(f"Failed to process file {file_name or 'Unknown'}: {str(e)}")
        return ""

def convert_docx_to_txt(file: UploadFile, file_store: FileStore, file_path: str) -> None:
    """
    将Word文档转换为txt文本文件并保存
    
    参数:
        file: 上传的Word文档文件
        file_store: 文件存储对象
        file_path: 文件保存路径
    """
    file.file.seek(0)
    docx_content = file.file.read()
    doc = Document(BytesIO(docx_content))

    # Extract text from the document
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)

    # Join the extracted text
    text_content = "\n".join(full_text)

    txt_file_path = docx_to_txt_filename(file_path)
    file_store.save_file(
        file_name=txt_file_path,
        content=BytesIO(text_content.encode("utf-8")),
        display_name=file.filename,
        file_origin=FileOrigin.CONNECTOR,
        file_type="text/plain",
    )

def docx_to_txt_filename(file_path: str) -> str:
    """
    将.docx文件路径转换为对应的.txt文件路径
    
    参数:
        file_path: docx文件路径
    返回:
        str: 对应的txt文件路径
    """
    return file_path.rsplit(".", 1)[0] + ".txt"
